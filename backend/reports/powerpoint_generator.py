from pathlib import Path

from pptx import Presentation


def generate_powerpoint_report(
    detections
):

    reports_dir = Path(
        "generated_reports"
    )

    reports_dir.mkdir(
        exist_ok=True
    )

    report_path = (

        reports_dir /

        "executive_briefing.pptx"
    )

    prs = Presentation()

    slide = prs.slides.add_slide(

        prs.slide_layouts[0]

    )

    slide.shapes.title.text = (

        "A3-ENDS Executive Briefing"
    )

    slide.placeholders[
        1
    ].text = (

        f"Incidents Analysed: "
        f"{len(detections)}"
    )

    prs.save(
        report_path
    )

    return str(
        report_path
    )