"""
Forensic Report Builders

Generates professional forensic reports in three formats:
  - DOCX  (python-docx)   — detailed narrative report
  - PPTX  (python-pptx)   — executive briefing slide deck
  - XLSX  (openpyxl)      — incident data spreadsheet

Each builder accepts a standardised `report_data` dict produced by
`build_report_data()` in report_generator.py.
"""

from __future__ import annotations
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List


# ---------------------------------------------------------------------------
# SHARED HELPERS
# ---------------------------------------------------------------------------

def _safe(value, fallback: str = "N/A") -> str:
    if value is None:
        return fallback
    return str(value)


# ---------------------------------------------------------------------------
# DOCX BUILDER
# ---------------------------------------------------------------------------

def build_docx(report_data: Dict[str, Any], output_path: Path) -> Path:
    """Build a Word forensic report."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # ── Title ──────────────────────────────────────────────────────────────
    title = doc.add_heading("A³-ENDS Forensic Incident Report", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.color.rgb = RGBColor(0, 0xA0, 0xFF)

    doc.add_paragraph(
        f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}  |  "
        f"Report ID: {report_data.get('report_id', 'N/A')}"
    ).alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()  # spacer

    # ── Executive Summary ──────────────────────────────────────────────────
    doc.add_heading("1. Executive Summary", level=1)
    llm_text: str = report_data.get("llm_analysis", "")
    if llm_text and not llm_text.startswith("{"):
        doc.add_paragraph(llm_text)
    else:
        sev    = _safe(report_data.get("severity"))
        attack = _safe(report_data.get("attack_type"))
        doc.add_paragraph(
            f"A {sev} severity {attack} attack was detected on "
            f"{_safe(report_data.get('source_ip'))} targeting "
            f"{_safe(report_data.get('dest_ip'))}. "
            f"Risk score: {_safe(report_data.get('risk_score'))}. "
            f"Analyst decision: {_safe(report_data.get('analyst_decision', 'pending'))}."
        )

    # ── Detection Details ──────────────────────────────────────────────────
    doc.add_heading("2. Detection Details", level=1)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Field"
    hdr_cells[1].text = "Value"
    for cell in hdr_cells:
        cell.paragraphs[0].runs[0].font.bold = True

    detail_rows = [
        ("Detection ID",   report_data.get("detection_id")),
        ("Timestamp",      report_data.get("timestamp")),
        ("Source IP",      report_data.get("source_ip")),
        ("Source Port",    report_data.get("source_port")),
        ("Destination IP", report_data.get("dest_ip")),
        ("Dest Port",      report_data.get("dest_port")),
        ("Protocol",       report_data.get("protocol")),
        ("Attack Type",    report_data.get("attack_type")),
        ("Severity",       report_data.get("severity")),
        ("Risk Score",     report_data.get("risk_score")),
        ("Confidence",     report_data.get("confidence")),
        ("Anomaly Score",  report_data.get("anomaly_score")),
    ]
    for label, value in detail_rows:
        row_cells = table.add_row().cells
        row_cells[0].text = label
        row_cells[1].text = _safe(value)

    doc.add_paragraph()

    # ── SHAP Explainability ────────────────────────────────────────────────
    doc.add_heading("3. SHAP Feature Contributions", level=1)
    shap: dict = report_data.get("shap_explanation") or {}
    if shap:
        shap_table = doc.add_table(rows=1, cols=2)
        shap_table.style = "Table Grid"
        h = shap_table.rows[0].cells
        h[0].text, h[1].text = "Feature", "SHAP Value"
        for cell in h:
            cell.paragraphs[0].runs[0].font.bold = True
        for feat, val in sorted(shap.items(), key=lambda x: abs(x[1]), reverse=True):
            r = shap_table.add_row().cells
            r[0].text = feat
            r[1].text = f"{val:+.6f}"
    else:
        doc.add_paragraph("SHAP explanation not available for this detection.")

    doc.add_paragraph()

    # ── Response Actions ───────────────────────────────────────────────────
    doc.add_heading("4. Response Actions", level=1)
    actions: list = report_data.get("response_actions") or []
    if actions:
        for action in actions:
            doc.add_paragraph(f"• {action}", style="List Bullet")
    else:
        doc.add_paragraph("No automated response actions were executed.")

    # ── Recommendations ────────────────────────────────────────────────────
    doc.add_heading("5. Recommendations", level=1)
    doc.add_paragraph(
        "1. Verify whether the source IP is an internal or external address.\n"
        "2. Check firewall and IDS logs for corroborating events.\n"
        "3. If the attack type is confirmed, escalate to the incident response team.\n"
        "4. Update threat intelligence feeds with observed IOCs."
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# PPTX BUILDER
# ---------------------------------------------------------------------------

def build_pptx(report_data: Dict[str, Any], output_path: Path) -> Path:
    """Build an executive PowerPoint briefing."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x06, 0x08, 0x0F)
    CYAN    = RGBColor(0x00, 0xD4, 0xFF)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    RED     = RGBColor(0xFF, 0x44, 0x44)
    AMBER   = RGBColor(0xFF, 0xA5, 0x00)

    SEV_COLOR = {
        "CRITICAL": RED,
        "HIGH":     RGBColor(0xFF, 0x66, 0x00),
        "MEDIUM":   AMBER,
        "LOW":      RGBColor(0x00, 0xC8, 0x53),
    }
    severity      = _safe(report_data.get("severity", "UNKNOWN"))
    attack_type   = _safe(report_data.get("attack_type", "Unknown"))
    sev_color     = SEV_COLOR.get(severity, WHITE)

    def _set_bg(slide):
        from pptx.util import Pt
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def _add_title(slide, text: str, sub: str = ""):
        tf = slide.shapes.title.text_frame if slide.shapes.title else None
        if tf is None:
            return
        tf.text = text
        tf.paragraphs[0].runs[0].font.color.rgb = CYAN
        tf.paragraphs[0].runs[0].font.size      = Pt(28)
        tf.paragraphs[0].runs[0].font.bold      = True
        if sub and len(tf.paragraphs) > 1:
            tf.paragraphs[1].text = sub
            tf.paragraphs[1].runs[0].font.color.rgb = WHITE
            tf.paragraphs[1].runs[0].font.size      = Pt(14)

    blank_layout = prs.slide_layouts[6]   # blank
    title_layout = prs.slide_layouts[0]   # title slide

    # ── SLIDE 1: Title ──────────────────────────────────────────────────
    s1 = prs.slides.add_slide(title_layout)
    _set_bg(s1)
    title_tf = s1.shapes.title.text_frame
    title_tf.text = "A³-ENDS Incident Briefing"
    title_tf.paragraphs[0].runs[0].font.color.rgb = CYAN
    title_tf.paragraphs[0].runs[0].font.size      = Pt(36)
    title_tf.paragraphs[0].runs[0].font.bold      = True
    if s1.placeholders[1]:
        s1.placeholders[1].text_frame.text = (
            f"{severity} Severity  |  {attack_type}  |  "
            f"{datetime.utcnow().strftime('%Y-%m-%d')}"
        )
        s1.placeholders[1].text_frame.paragraphs[0].runs[0].font.color.rgb = sev_color

    # ── SLIDE 2: Threat Overview ────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    _set_bg(s2)

    def _add_label(slide, x, y, w, h, text, color=WHITE, size=14, bold=False):
        txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text            = text
        run.font.color.rgb  = color
        run.font.size       = Pt(size)
        run.font.bold       = bold

    _add_label(s2, 0.5, 0.3, 12, 0.7, "Threat Overview", CYAN, 24, bold=True)

    kpis = [
        ("Attack Type",   attack_type,                                     sev_color),
        ("Severity",      severity,                                         sev_color),
        ("Source IP",     _safe(report_data.get("source_ip")),             WHITE),
        ("Target IP",     _safe(report_data.get("dest_ip")),              WHITE),
        ("Risk Score",    _safe(report_data.get("risk_score")),            AMBER),
        ("Confidence",    _safe(report_data.get("confidence")),            CYAN),
    ]
    for i, (label, value, color) in enumerate(kpis):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.3
        y = 1.2 + row * 2.0
        _add_label(s2, x, y,       4.0, 0.4, label, RGBColor(0x5A, 0x9A, 0xBB), 11)
        _add_label(s2, x, y + 0.4, 4.0, 0.8, value, color, 18, bold=True)

    # ── SLIDE 3: SHAP Explanation ────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    _set_bg(s3)
    _add_label(s3, 0.5, 0.3, 12, 0.7, "AI Explainability — Top Feature Drivers", CYAN, 24, bold=True)

    shap: dict = report_data.get("shap_explanation") or {}
    if shap:
        sorted_shap = sorted(shap.items(), key=lambda x: abs(x[1]), reverse=True)[:6]
        max_abs = max(abs(v) for _, v in sorted_shap) if sorted_shap else 1.0
        for i, (feat, val) in enumerate(sorted_shap):
            y  = 1.2 + i * 0.85
            clr = RGBColor(0xFF, 0x44, 0x44) if val >= 0 else RGBColor(0x5A, 0x9A, 0xBB)
            bar_w = max(0.1, 7.0 * abs(val) / max_abs)
            _add_label(s3, 0.5, y, 2.5, 0.5, feat, WHITE, 12)
            bar = s3.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(3.2), Inches(y + 0.1), Inches(bar_w), Inches(0.3)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = clr
            bar.line.fill.background()
            _add_label(s3, 3.2 + bar_w + 0.1, y, 2.0, 0.5, f"{val:+.4f}", clr, 11)
    else:
        _add_label(s3, 0.5, 1.5, 12, 1.0, "SHAP explanation not available.", WHITE, 14)

    # ── SLIDE 4: Recommended Actions ────────────────────────────────────
    s4 = prs.slides.add_slide(blank_layout)
    _set_bg(s4)
    _add_label(s4, 0.5, 0.3, 12, 0.7, "Recommended Response Actions", CYAN, 24, bold=True)

    actions: list = report_data.get("response_actions") or [
        "Review and confirm the detection",
        "Check firewall and access logs",
        "Escalate to incident response if confirmed",
    ]
    for i, action in enumerate(actions[:8]):
        _add_label(s4, 0.8, 1.2 + i * 0.75, 12.0, 0.6, f"▸  {action}", WHITE, 14)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# XLSX BUILDER
# ---------------------------------------------------------------------------

def build_xlsx(report_data: Dict[str, Any], output_path: Path) -> Path:
    """Build an Excel incident data spreadsheet."""
    from openpyxl import Workbook
    from openpyxl.styles import (
        Font, PatternFill, Alignment, Border, Side, numbers
    )

    DARK   = "FF06080F"
    CYAN   = "FF00D4FF"
    HEADER = "FF0A1020"
    WHITE  = "FFFFFFFF"
    AMBER  = "FFFFA500"
    RED    = "FFFF4444"
    GREEN  = "FF00C853"

    wb = Workbook()

    # ── Sheet 1: Incident Summary ────────────────────────────────────────
    ws1 = wb.active
    ws1.title = "Incident Summary"
    ws1.sheet_properties.tabColor = "00D4FF"

    def _hdr_cell(ws, row, col, text, width=20):
        cell = ws.cell(row=row, column=col, value=text)
        cell.font      = Font(bold=True, color=CYAN, name="Consolas", size=10)
        cell.fill      = PatternFill(fill_type="solid", fgColor=HEADER)
        cell.alignment = Alignment(horizontal="left", vertical="center")
        ws.column_dimensions[cell.column_letter].width = width
        return cell

    def _data_cell(ws, row, col, text, color=WHITE, bold=False):
        cell = ws.cell(row=row, column=col, value=text)
        cell.font      = Font(color=color, name="Consolas", size=10, bold=bold)
        cell.fill      = PatternFill(fill_type="solid", fgColor=DARK)
        cell.alignment = Alignment(horizontal="left", vertical="center")
        return cell

    # Title row
    ws1.row_dimensions[1].height = 28
    t = ws1.cell(row=1, column=1, value="A³-ENDS FORENSIC INCIDENT REPORT")
    t.font      = Font(bold=True, color=CYAN, name="Consolas", size=14)
    t.fill      = PatternFill(fill_type="solid", fgColor=DARK)
    t.alignment = Alignment(horizontal="center", vertical="center")
    ws1.merge_cells("A1:D1")

    headers = ["Field", "Value", "Notes"]
    for i, h in enumerate(headers, 1):
        _hdr_cell(ws1, 2, i, h, width=25)

    summary_rows = [
        ("Report ID",      report_data.get("report_id"),       ""),
        ("Generated At",   datetime.utcnow().isoformat(),       "UTC"),
        ("Detection ID",   report_data.get("detection_id"),    ""),
        ("Timestamp",      report_data.get("timestamp"),        "UTC"),
        ("Attack Type",    report_data.get("attack_type"),      "8-class LightGBM"),
        ("Severity",       report_data.get("severity"),         "LOW/MEDIUM/HIGH/CRITICAL"),
        ("Risk Score",     report_data.get("risk_score"),       "0–100"),
        ("Confidence",     report_data.get("confidence"),       "0.0–1.0"),
        ("Anomaly Score",  report_data.get("anomaly_score"),    "MSE reconstruction error"),
        ("Source IP",      report_data.get("source_ip"),        ""),
        ("Source Port",    report_data.get("source_port"),      ""),
        ("Destination IP", report_data.get("dest_ip"),          ""),
        ("Dest Port",      report_data.get("dest_port"),        ""),
        ("Protocol",       report_data.get("protocol"),         ""),
        ("Analyst Decision", report_data.get("analyst_decision", "pending"), ""),
        ("Analyst Notes",  report_data.get("analyst_notes"),   ""),
    ]
    for r, (field, value, note) in enumerate(summary_rows, start=3):
        sev = report_data.get("severity", "")
        val_color = RED if field == "Severity" and sev in ("CRITICAL", "HIGH") else WHITE
        _data_cell(ws1, r, 1, field)
        _data_cell(ws1, r, 2, _safe(value), color=val_color)
        _data_cell(ws1, r, 3, note)
        ws1.row_dimensions[r].height = 18

    # ── Sheet 2: SHAP Explanation ────────────────────────────────────────
    ws2 = wb.create_sheet("SHAP Explanation")
    ws2.sheet_properties.tabColor = "00D4FF"

    for col, (h, w) in enumerate([("Feature", 22), ("SHAP Value", 18), ("Abs Importance", 18), ("Direction", 14)], 1):
        _hdr_cell(ws2, 1, col, h, width=w)

    shap: dict = report_data.get("shap_explanation") or {}
    sorted_shap = sorted(shap.items(), key=lambda x: abs(x[1]), reverse=True) if shap else []
    for r, (feat, val) in enumerate(sorted_shap, start=2):
        _data_cell(ws2, r, 1, feat)
        _data_cell(ws2, r, 2, round(float(val), 6), color=RED if val >= 0 else GREEN)
        _data_cell(ws2, r, 3, round(abs(float(val)), 6))
        _data_cell(ws2, r, 4, "↑ Positive" if val >= 0 else "↓ Negative",
                   color=RED if val >= 0 else GREEN)
        ws2.row_dimensions[r].height = 18

    # ── Sheet 3: Response Actions ─────────────────────────────────────────
    ws3 = wb.create_sheet("Response Actions")
    ws3.sheet_properties.tabColor = "FFA500"

    _hdr_cell(ws3, 1, 1, "Action", width=60)
    _hdr_cell(ws3, 1, 2, "Status", width=18)

    actions: list = report_data.get("response_actions") or []
    for r, action in enumerate(actions, start=2):
        _data_cell(ws3, r, 1, action)
        _data_cell(ws3, r, 2, "Executed", color=GREEN)
        ws3.row_dimensions[r].height = 18

    # ── Sheet 4: LLM Analysis ─────────────────────────────────────────────
    ws4 = wb.create_sheet("LLM Analysis")
    ws4.sheet_properties.tabColor = "5A9ABB"

    _hdr_cell(ws4, 1, 1, "LLM Forensic Analysis", width=100)
    llm_text = _safe(report_data.get("llm_analysis", "Not generated."))
    for r, line in enumerate(llm_text.splitlines(), start=2):
        c = ws4.cell(row=r, column=1, value=line)
        c.font      = Font(color=WHITE, name="Consolas", size=10)
        c.fill      = PatternFill(fill_type="solid", fgColor=DARK)
        c.alignment = Alignment(wrap_text=True)
        ws4.row_dimensions[r].height = 15

    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output_path))
    return output_path
